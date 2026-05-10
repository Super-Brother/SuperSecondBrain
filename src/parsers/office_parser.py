"""Office 文档解析器（Word / PPT / Excel）"""

import os
import re
from pathlib import Path

from src.parsers.base_parser import BaseParser, Document


class WordParser(BaseParser):
    """Word 文档解析，依赖 python-docx"""

    def __init__(self, base_path: str = ""):
        self.base_path = Path(base_path).resolve() if base_path else Path.cwd()

    def parse_file(self, file_path: str) -> Document:
        path = Path(file_path).resolve()

        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        content = "\n\n".join(paragraphs)
        title = path.stem

        if doc.core_properties.title:
            title = doc.core_properties.title

        rel_path = str(path.relative_to(self.base_path)) if path.is_relative_to(self.base_path) else path.name
        folder = str(path.parent.relative_to(self.base_path)) if path.parent != self.base_path else "root"

        return Document(
            title=title,
            content=content,
            source_file=str(path),
            relative_path=rel_path,
            folder=folder,
            content_hash=self.compute_hash(content),
            metadata={"page_count": len(doc.sections), "paragraph_count": len(paragraphs)},
        )

    def parse_directory(self, dir_path: str, exclude_dirs: list[str] = None) -> list[Document]:
        if exclude_dirs is None:
            exclude_dirs = [".git"]

        docs = []
        for docx_file in Path(dir_path).rglob("*.docx"):
            if any(part in exclude_dirs for part in docx_file.parts):
                continue
            try:
                doc = self.parse_file(str(docx_file))
                if len(doc.content) > 50:
                    docs.append(doc)
            except Exception as e:
                print(f"[WARN] Word解析失败 {docx_file}: {e}")

        return docs


class PPTParser(BaseParser):
    """PPT 文档解析，依赖 python-pptx"""

    def __init__(self, base_path: str = ""):
        self.base_path = Path(base_path).resolve() if base_path else Path.cwd()

    def parse_file(self, file_path: str) -> Document:
        path = Path(file_path).resolve()

        from pptx import Presentation

        prs = Presentation(str(path))

        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides_text.append(f"--- 第{i}页 ---\n" + "\n".join(slide_texts))

        content = "\n\n".join(slides_text)
        title = path.stem

        rel_path = str(path.relative_to(self.base_path)) if path.is_relative_to(self.base_path) else path.name
        folder = str(path.parent.relative_to(self.base_path)) if path.parent != self.base_path else "root"

        return Document(
            title=title,
            content=content,
            source_file=str(path),
            relative_path=rel_path,
            folder=folder,
            content_hash=self.compute_hash(content),
            metadata={"slide_count": len(prs.slides)},
        )

    def parse_directory(self, dir_path: str, exclude_dirs: list[str] = None) -> list[Document]:
        if exclude_dirs is None:
            exclude_dirs = [".git"]

        docs = []
        for ppt_file in Path(dir_path).rglob("*.pptx"):
            if any(part in exclude_dirs for part in ppt_file.parts):
                continue
            try:
                doc = self.parse_file(str(ppt_file))
                if len(doc.content) > 50:
                    docs.append(doc)
            except Exception as e:
                print(f"[WARN] PPT解析失败 {ppt_file}: {e}")

        return docs


class ExcelParser(BaseParser):
    """Excel 文档解析，依赖 openpyxl"""

    def __init__(self, base_path: str = ""):
        self.base_path = Path(base_path).resolve() if base_path else Path.cwd()

    def parse_file(self, file_path: str) -> Document:
        path = Path(file_path).resolve()

        from openpyxl import load_workbook

        wb = load_workbook(str(path), data_only=True)

        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets_text.append(f"--- 工作表: {sheet_name} ---\n" + "\n".join(rows))

        content = "\n\n".join(sheets_text)
        title = path.stem

        rel_path = str(path.relative_to(self.base_path)) if path.is_relative_to(self.base_path) else path.name
        folder = str(path.parent.relative_to(self.base_path)) if path.parent != self.base_path else "root"

        return Document(
            title=title,
            content=content,
            source_file=str(path),
            relative_path=rel_path,
            folder=folder,
            content_hash=self.compute_hash(content),
            metadata={"sheet_count": len(wb.sheetnames)},
        )

    def parse_directory(self, dir_path: str, exclude_dirs: list[str] = None) -> list[Document]:
        if exclude_dirs is None:
            exclude_dirs = [".git"]

        docs = []
        for xlsx_file in Path(dir_path).rglob("*.xlsx"):
            if any(part in exclude_dirs for part in xlsx_file.parts):
                continue
            try:
                doc = self.parse_file(str(xlsx_file))
                if len(doc.content) > 50:
                    docs.append(doc)
            except Exception as e:
                print(f"[WARN] Excel解析失败 {xlsx_file}: {e}")

        return docs
